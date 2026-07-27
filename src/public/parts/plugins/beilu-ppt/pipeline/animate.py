# -*- coding: utf-8 -*-
"""
animate.py —— v0.3 动画后处理（契约 v0.3 / 探针实测结构）

线路: python-pptx 打开 pptx → 每页取真实 shape_id（禁 regex 抓 cNvPr）
      → lxml etree 建 <p:timing> 子树（禁字符串拼嵌套XML）→ 追加到 <p:sld> 尾
      → 保存另存/原地。

为什么这么做（血泪根因，全部主AI探针实测）:
  - spid 必须来自真实 shape: python-pptx 的 sh.shape_id 才是入场目标; 若 regex 抓 <p:cNvPr>
    会把 spTree 根 id=1（nvGrpSpPr，不是 shape）也算进去 → PowerPoint 挂出一个空 shape 幽灵效果。
  - lxml etree 建树而非字符串拼接: 探针字符串版曾少闭合一个根 <p:par> 直接坏文件。
  - 结构层级 (探针验证):
      timing > tnLst > par > cTn(id=1 tmRoot) > childTnLst
        > seq(mainSeq) > cTn(id=2 mainSeq) > childTnLst
          > 每 shape 一个 clickEffect: par>cTn(100/110/..)>childTnLst
              > par>cTn(101/..)>childTnLst
                > par>cTn(102 presetID nodeType=clickEffect)>childTnLst
                    > set(style.visibility→visible)  + animEffect(transition=in filter=fade|wipe)
        seq 尾带 prevCondLst/nextCondLst(带 sldTgt onPrev/onNext)
  - cTn id 全文件唯一递增（探针 100,101,102,103,104 / 110,... 每 shape 一簇）。
  - 幂等: slide XML 已有 <p:timing> 则跳过该页。

模式:
  fade: presetID=10 presetClass=entr filter="fade"
  wipe: presetID=22 presetClass=entr filter="wipe(fromBottom)"
"""
from lxml import etree
from pptx import Presentation

# lxml 命名空间（契约指定 p= presentationml）
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
NSMAP = {"p": P_NS}


def _q(tag):
    """p:tag → Clark notation 全限定名，供 lxml 建元素。"""
    return "{%s}%s" % (P_NS, tag)


def _sub(parent, tag, **attrs):
    """建子元素并挂到 parent；属性无命名空间（与探针 XML 一致）。"""
    el = etree.SubElement(parent, _q(tag))
    for k, v in attrs.items():
        el.set(k, str(v))
    return el


# 模式 → (presetID, filter)。v2.3 扩域: 全部走同一 animEffect transition=in + filter 机制
# （与 fade/wipe 同结构, PowerPoint 标准 entr preset + 滤镜名, 不引入新 XML 形状）。
_MODES = {
    "fade": (10, "fade"),
    "wipe": (22, "wipe(fromBottom)"),
    "blinds": (3, "blinds(horizontal)"),
    "box": (4, "box(in)"),
    "checkerboard": (5, "checkerboard(across)"),
    "circle": (6, "circle(out)"),
    "dissolve": (9, "dissolve"),
}
MODES = _MODES  # 对外单源（pipeline 值域校验用, 防 JS/python 双源枚举）


def _build_click_effect(parent_childTnLst, base_id, spid, preset_id, flt):
    """
    为单个 shape 建一簇 clickEffect（探针三层 par 嵌套 + set + animEffect）。
    base_id: 该簇的起始 cTn id（用 100/110/120...，簇内 +0..+4）。
    返回下一个可用 base_id（本簇消耗 5 个 id）。
    """
    # 第一层 par: cTn(base+0) fill=hold, stCondLst cond delay=indefinite
    par0 = _sub(parent_childTnLst, "par")
    cTn0 = _sub(par0, "cTn", id=base_id, fill="hold")
    st0 = _sub(cTn0, "stCondLst")
    _sub(st0, "cond", delay="indefinite")
    ctl0 = _sub(cTn0, "childTnLst")

    # 第二层 par: cTn(base+1) fill=hold, cond delay=0
    par1 = _sub(ctl0, "par")
    cTn1 = _sub(par1, "cTn", id=base_id + 1, fill="hold")
    st1 = _sub(cTn1, "stCondLst")
    _sub(st1, "cond", delay="0")
    ctl1 = _sub(cTn1, "childTnLst")

    # 第三层 par: cTn(base+2) presetID/presetClass nodeType=clickEffect
    par2 = _sub(ctl1, "par")
    cTn2 = _sub(par2, "cTn", id=base_id + 2, presetID=preset_id,
                presetClass="entr", presetSubtype="0", fill="hold",
                nodeType="clickEffect")
    st2 = _sub(cTn2, "stCondLst")
    _sub(st2, "cond", delay="0")
    ctl2 = _sub(cTn2, "childTnLst")

    # set: style.visibility → visible （cTn id=base+3）
    setel = _sub(ctl2, "set")
    cbhvr_s = _sub(setel, "cBhvr")
    cTn3 = _sub(cbhvr_s, "cTn", id=base_id + 3, dur="1", fill="hold")
    st3 = _sub(cTn3, "stCondLst")
    _sub(st3, "cond", delay="0")
    tgt_s = _sub(cbhvr_s, "tgtEl")
    _sub(tgt_s, "spTgt", spid=spid)
    anl = _sub(cbhvr_s, "attrNameLst")
    anm = _sub(anl, "attrName")
    anm.text = "style.visibility"
    toel = _sub(setel, "to")
    _sub(toel, "strVal", val="visible")

    # animEffect: transition=in filter=fade|wipe （cTn id=base+4）
    ae = _sub(ctl2, "animEffect", transition="in", filter=flt)
    cbhvr_a = _sub(ae, "cBhvr")
    _sub(cbhvr_a, "cTn", id=base_id + 4, dur="500")
    tgt_a = _sub(cbhvr_a, "tgtEl")
    _sub(tgt_a, "spTgt", spid=spid)

    return base_id + 10  # 每簇留 10 号段，簇间不撞


def _build_timing(spids, preset_id, flt):
    """建整棵 <p:timing> lxml 元素。spids: 该页有序 shape_id 列表（title 已排首）。"""
    timing = etree.Element(_q("timing"), nsmap=NSMAP)
    tnLst = _sub(timing, "tnLst")
    par_root = _sub(tnLst, "par")
    cTn_root = _sub(par_root, "cTn", id="1", dur="indefinite",
                    restart="never", nodeType="tmRoot")
    ctl_root = _sub(cTn_root, "childTnLst")

    seq = _sub(ctl_root, "seq", concurrent="1", nextAc="seek")
    cTn_seq = _sub(seq, "cTn", id="2", dur="indefinite", nodeType="mainSeq")
    ctl_seq = _sub(cTn_seq, "childTnLst")

    base = 100
    for spid in spids:
        base = _build_click_effect(ctl_seq, base, str(spid), preset_id, flt)

    # seq 尾: prevCondLst / nextCondLst 带 sldTgt
    prev = _sub(seq, "prevCondLst")
    condp = _sub(prev, "cond", evt="onPrev", delay="0")
    tgtp = _sub(condp, "tgtEl")
    _sub(tgtp, "sldTgt")
    nxt = _sub(seq, "nextCondLst")
    condn = _sub(nxt, "cond", evt="onNext", delay="0")
    tgtn = _sub(condn, "tgtEl")
    _sub(tgtn, "sldTgt")

    return timing


def _ordered_spids(slide, slide_title):
    """
    返回该页有序 shape_id 列表，title shape 排第一。
    title 判定: 文本以标题内容开头的第一个 shape; 无匹配则取 top(y) 最小者。
    其余按文档顺序（shapes 遍历序）。
    """
    # v0.6: 跳过全页背景图(覆盖≥90%画布的 Picture)——背景不该等点击才淡入
    part = slide.part.package.presentation_part.presentation
    sw, sh_h = part.slide_width, part.slide_height

    def _is_fullpage_bg(s):
        try:
            return (s.shape_type == 13  # PICTURE
                    and s.width is not None and s.height is not None
                    and s.width >= sw * 0.9 and s.height >= sh_h * 0.9)
        except Exception:
            return False

    def _is_decor(s):
        """v0.9: 纯装饰形状(卡片底/竖条=无文字AutoShape)不挂动画——
        否则和其上文字分离出现且点击次数翻倍(实测39效果/6页)。图片/图表保留。"""
        try:
            if s.shape_type == 13 or getattr(s, "has_chart", False):  # Picture/Chart
                return False
            if s.has_text_frame and s.text_frame.text.strip():
                return False
            return True
        except Exception:
            return False

    shapes = [s for s in slide.shapes if not _is_fullpage_bg(s) and not _is_decor(s)]
    if not shapes:
        return []

    title_sh = None
    if slide_title:
        t = slide_title.strip()
        for sh in shapes:
            try:
                if sh.has_text_frame and sh.text_frame.text.strip().startswith(t):
                    title_sh = sh
                    break
            except Exception:
                pass
    if title_sh is None:
        # 按 y 最小（top 可能 None，兜底大值）
        title_sh = min(shapes, key=lambda s: (s.top if s.top is not None else 1 << 60))

    ordered = [title_sh] + [s for s in shapes if s is not title_sh]
    return [s.shape_id for s in ordered]


def _has_timing(slide):
    """幂等检查: slide XML 已含 <p:timing> 则 True。"""
    return slide._element.find(_q("timing")) is not None


def add_animations(pptx_path, mode="fade", out_path=None, per_slide=None):
    """
    给 pptx 每页每 shape 注入 clickEffect 入场动画。
    mode: MODES 键（deck 级缺省）。out_path: None=原地保存, 否则另存。
    per_slide: 可选 [mode|None, ...] 按页序覆盖（v2.3, spec 的 slide.animate;
    None=用 deck 级; "none"=该页不加动画）。
    返回最终 pptx 路径（str）。
    """
    if mode not in _MODES:
        raise ValueError("mode 必须是 %s，收到 %r" % (list(_MODES), mode))

    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        if _has_timing(slide):
            continue  # 幂等
        sl_mode = mode
        if per_slide and i < len(per_slide) and per_slide[i]:
            sl_mode = str(per_slide[i]).lower()
            if sl_mode == "none":
                continue  # 页级显式关
            if sl_mode not in _MODES:
                sl_mode = mode  # 未知页级值退 deck 级（值域校验在 pipeline 已回喂信号）
        preset_id, flt = _MODES[sl_mode]
        title = ""
        try:
            if slide.shapes.title is not None:
                title = slide.shapes.title.text or ""
        except Exception:
            title = ""
        # 空白版式无占位 title，用文本/最小y启发式
        spids = _ordered_spids(slide, title)
        if not spids:
            continue
        timing = _build_timing(spids, preset_id, flt)
        slide._element.append(timing)  # <p:timing> 作 <p:sld> 尾子元素

    final = out_path or pptx_path
    prs.save(final)
    return str(final)


# ---------------- 自测 ----------------

if __name__ == "__main__":
    import sys
    from pathlib import Path
    from xml.dom import minidom
    import zipfile

    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    HERE = Path(__file__).parent
    src = HERE / "out" / "deck.pptx"
    dst = HERE / "out" / "deck_fade.pptx"
    assert src.exists(), "缺 out/deck.pptx，先跑 pipeline.py"

    out = add_animations(str(src), mode="fade", out_path=str(dst))
    print("[生成]", out)

    # 验1: python-pptx 重开 OK + 收集每页 shape 数
    prs = Presentation(out)
    shape_counts = [len(list(s.shapes)) for s in prs.slides]
    print("[验1] python-pptx 重开 OK, 每页 shape 数:", shape_counts)

    # 验2 & 验3: 每页 slide XML minidom 良构 + timing 存在 + spTgt 数 == shape 数
    z = zipfile.ZipFile(out)
    slide_names = sorted(n for n in z.namelist()
                         if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
    for i, name in enumerate(slide_names):
        raw = z.read(name)
        dom = minidom.parseString(raw)  # 良构则不抛
        xml = raw.decode("utf-8")
        assert "<p:timing" in xml, "%s 缺 timing" % name
        # 每 shape 一个 set + 一个 animEffect = 2 个 spTgt；契约要 spTgt 数 == shape 数,
        # 这里核入场效果数（唯一 spid 数 == shape 数）更贴切: 用 set 的 spTgt 计数
        n_set_sptgt = xml.count("<p:set>")
        n_shapes = shape_counts[i]
        assert n_set_sptgt == n_shapes, \
            "%s 效果数 %d != shape 数 %d" % (name, n_set_sptgt, n_shapes)
        n_sptgt = xml.count("<p:spTgt")
        print("[验2/3] %s 良构OK, 入场效果=%d shape=%d spTgt总=%d(2/shape)"
              % (name, n_set_sptgt, n_shapes, n_sptgt))

    # 验4: 幂等 —— 再跑一次不应改变 timing 数
    add_animations(out, mode="fade")  # 原地, 应全跳过
    z2 = zipfile.ZipFile(out)
    for name in slide_names:
        cnt = z2.read(name).decode("utf-8").count("<p:timing")
        assert cnt == 1, "%s timing 数 %d != 1（幂等破了）" % (name, cnt)
    print("[验4] 幂等 OK: 重跑后每页仍只 1 个 timing")

    print("[自测全过]")

# -*- coding: utf-8 -*-
"""
pptx_paint.py —— v3 薄画笔: 场景图 → .pptx（零测量零视觉决策）

契约（v3框架设计 §三/§五）: 只认 8 种原语, 按 z 升序照画; 所有颜色/坐标/runs 已由
compose 决策完毕。产物全原生对象（textbox/shape/connector/table/chart/picture）=
每个元素都是组件, 全可编辑（凛倾 0717 组装铁律）。
preview_only 元素跳过（溢出标记等预览诊断不进交付件）。
单元素失败隔离: 记 errors 继续画（容错矩阵 paint_error）。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

import tokens

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_CHART = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "line": XL_CHART_TYPE.LINE,
          "pie": XL_CHART_TYPE.PIE}


def _rgb(h):
    v = tokens.hex_rgb(h)
    return RGBColor(*v) if v else None


def _set_alpha(shp, alpha):
    """形状填充透明度: OOXML a:alpha 注入（python-pptx 无 API; val=不透明度%×1000）。"""
    try:
        srgb = shp._element.spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        srgb.append(srgb.makeelement(qn("a:alpha"),
                                     {"val": str(int(max(0.0, min(1.0, alpha)) * 100000))}))
    except Exception:
        pass  # 透明度是装饰增强, 失败退不透明


def _paint_text(slide, el, font_name):
    tb = slide.shapes.add_textbox(Inches(el["x"]), Inches(el["y"]),
                                  Inches(el["w"]), Inches(el["h"]))
    tf = tb.text_frame
    tf.word_wrap = bool(el.get("wrap", True))
    tf.auto_size = MSO_AUTO_SIZE.NONE
    if el.get("anchor") == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    align = _ALIGN.get(el.get("align", "left"), PP_ALIGN.LEFT)
    for i, row in enumerate(el.get("runs") or [[]]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for r in row:
            run = p.add_run()
            run.text = r.get("t", "")
            run.font.name = font_name
            run.font.size = Pt(r.get("pt", 14))
            run.font.bold = bool(r.get("bold"))
            c = _rgb(r.get("color"))
            if c is not None:
                run.font.color.rgb = c


def _paint_shape(slide, el):
    kind = MSO_SHAPE.OVAL if el.get("shape") == "oval" else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(el["x"]), Inches(el["y"]),
                                 Inches(el["w"]), Inches(el["h"]))
    if el.get("fill"):
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(el["fill"])
        if el.get("alpha") is not None:
            _set_alpha(shp, el["alpha"])
    else:
        shp.fill.background()
    if el.get("line_color"):
        shp.line.color.rgb = _rgb(el["line_color"])
        shp.line.width = Pt(el.get("line_pt", 1.0))
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False


def _paint_line(slide, el):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(el["x1"]), Inches(el["y1"]),
        Inches(el["x2"]), Inches(el["y2"]))
    conn.line.color.rgb = _rgb(el.get("color")) or RGBColor(0, 0, 0)
    conn.line.width = Pt(el.get("width_pt", 1.0))
    conn.shadow.inherit = False
    if el.get("arrow"):
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))


def _paint_gradient(slide, el):
    """渐变底: 全幅矩形+双 stop 渐变（API 缺失时退纯色 from——诚实降级不空底）。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(el["x"]), Inches(el["y"]),
                                 Inches(el["w"]), Inches(el["h"]))
    shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.fill.gradient()
        stops = shp.fill.gradient_stops
        stops[0].color.rgb = _rgb(el.get("from")) or RGBColor(0, 0, 0)
        stops[1].color.rgb = _rgb(el.get("to")) or RGBColor(0x33, 0x33, 0x33)
        try:
            shp.fill.gradient_angle = float(el.get("angle", 135)) % 360
        except Exception:
            pass  # 角度设不上保持默认对角
    except Exception:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(el.get("from")) or RGBColor(0x20, 0x20, 0x20)


def _paint_image(slide, el):
    slide.shapes.add_picture(el["src"], Inches(el["x"]), Inches(el["y"]),
                             Inches(el["w"]), Inches(el["h"]))


def _paint_table(slide, el, font_name):
    cells = el.get("cells") or []
    nrows = len(cells)
    ncols = max((len(r) for r in cells), default=1)
    if not nrows or not ncols:
        return
    row_h = el.get("row_h_in") or []
    gframe = slide.shapes.add_table(
        nrows, ncols, Inches(el["x"]), Inches(el["y"]),
        Inches(el["w"]), Inches(sum(row_h) if row_h else 0.4 * nrows))
    tbl = gframe.table
    tbl.first_row = True
    tbl.horz_banding = True
    for j, cw in enumerate((el.get("col_w_in") or [])[:ncols]):
        tbl.columns[j].width = Inches(cw)
    for i, rh in enumerate(row_h[:nrows]):
        tbl.rows[i].height = Inches(rh)
    m = el.get("margin", 0.06)
    for i, row in enumerate(cells):
        for j in range(ncols):
            cinfo = row[j] if j < len(row) else {"full": [], "fill": None}
            cell = tbl.cell(i, j)
            if cinfo.get("fill"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(cinfo["fill"])
            cell.margin_left = cell.margin_right = Inches(m)
            cell.margin_top = cell.margin_bottom = Inches(m / 2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            # 契约 v1.5 不变: 单元格喂原文（full 单行 runs）, PowerPoint 原生 wrap
            for r in (cinfo.get("full") or []):
                run = p.add_run()
                run.text = r.get("t", "")
                run.font.name = font_name
                run.font.size = Pt(r.get("pt", 12))
                run.font.bold = bool(r.get("bold"))
                c = _rgb(r.get("color"))
                if c is not None:
                    run.font.color.rgb = c


def _paint_chart(slide, el):
    data = CategoryChartData()
    data.categories = el.get("categories") or [""]
    for s in (el.get("series") or []):
        vals = []
        for v in (s.get("values") or []):
            try:
                vals.append(float(v))
            except (TypeError, ValueError):
                vals.append(0.0)
        data.add_series(s.get("name", ""), vals)
    gframe = slide.shapes.add_chart(
        _CHART.get(el.get("chart_type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED),
        Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"]), data)
    chart = gframe.chart
    if el.get("title"):
        chart.has_title = True
        chart.chart_title.text_frame.text = el["title"]
    else:
        chart.has_title = False
    if len(el.get("series") or []) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False


_PAINTERS = {"text": None, "shape": None, "line": None, "gradient": None,
             "image": None, "table": None, "chart": None, "group": None}


def render_pptx_scene(sc, out_path):
    """场景图 → .pptx。返回 (out_path, errors)。"""
    prs = Presentation()
    prs.slide_width = Inches(sc["slide_w_in"])
    prs.slide_height = Inches(sc["slide_h_in"])
    blank = prs.slide_layouts[6]
    font_name = sc.get("font", "微软雅黑")
    errors = []

    def paint(slide, el):
        p = el.get("p")
        if p == "text":
            _paint_text(slide, el, font_name)
        elif p == "shape":
            _paint_shape(slide, el)
        elif p == "line":
            _paint_line(slide, el)
        elif p == "gradient":
            _paint_gradient(slide, el)
        elif p == "image":
            _paint_image(slide, el)
        elif p == "table":
            _paint_table(slide, el, font_name)
        elif p == "chart":
            _paint_chart(slide, el)
        elif p == "group":
            for sub in (el.get("items") or []):
                paint(slide, sub)

    for sl in sc.get("slides", []):
        slide = prs.slides.add_slide(blank)
        for el in sl.get("elements", []):
            if el.get("preview_only"):
                continue
            try:
                paint(slide, el)
            except Exception as e:
                errors.append({"type": "paint_error", "slide": sl.get("id"),
                               "p": el.get("p"), "detail": str(e)[:120]})
    # v3.2 组装容错: ①临时文件+原子替换（崩溃不留半写坏件）②目标被占用（用户在
    # PowerPoint 里开着 deck.pptx 再让 AI 改=高频真实场景, Windows 文件锁 PermissionError）
    # → 另存时间戳副本+信号, 不炸管线
    import os
    import time
    tmp = out_path + ".tmp"
    prs.save(tmp)
    try:
        os.replace(tmp, out_path)
    except PermissionError:
        alt = out_path[:-5] + "_%s.pptx" % time.strftime("%H%M%S")
        os.replace(tmp, alt)
        errors.append({"type": "pptx_locked", "detail":
                       "deck.pptx 被占用(可能正在 PowerPoint 中打开), 本次已另存 %s——"
                       "关闭文件后重生成可恢复同名覆盖" % os.path.basename(alt)})
        return alt, errors
    return out_path, errors
